# masters/utils/import_export_tools.py
"""
Outil centralisé d'Import-Export pour la plateforme ESFé Master.
Compatible : Excel, CSV, JSON, PDF, Word, PowerPoint.
Fonctions génériques pour tous les modèles (étudiants, modules, enseignants, examens...).
Auteur : Mohamed Aly Camara x ChatGPT - 2025
"""

import io, pandas as pd, json
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple, Type
from django.apps import apps
from django.db import transaction
from django.http import HttpResponse
from django.core.exceptions import ValidationError
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from docx import Document


# ============================================================
# ⚙️ 1️⃣ OUTILS GÉNÉRAUX : normalisation / massage / validation
# ============================================================
def normalize_str(value: Any) -> str:
    """Nettoie les chaînes (espaces, casse, caractères spéciaux)."""
    if not value:
        return ""
    return str(value).strip().replace("\n", " ").replace("\r", "")

def normalize_number(value: Any) -> Optional[float]:
    """Convertit proprement une valeur numérique ou vide."""
    try:
        if value in ("", None, "NaN"):
            return None
        return float(value)
    except Exception:
        return None

def normalize_date(value: Any) -> Optional[datetime]:
    """Convertit proprement une date ou renvoie None."""
    if not value:
        return None
    if isinstance(value, datetime):
        return value
    try:
        return pd.to_datetime(value, dayfirst=True, errors="coerce")
    except Exception:
        return None

def massage_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """
    Applique une normalisation de base à un DataFrame :
      - Trim des chaînes
      - Conversion des NaN
      - Nettoyage des colonnes
    """
    df.columns = [normalize_str(c).lower().strip() for c in df.columns]
    for col in df.columns:
        if df[col].dtype == "object":
            df[col] = df[col].fillna("").astype(str).str.strip()
        else:
            df[col] = df[col].fillna(0)
    return df


# ============================================================
# 📥 2️⃣ IMPORT GÉNÉRIQUE : Excel / CSV / JSON vers Base Django
# ============================================================
@transaction.atomic
def import_from_file(model_name: str, file, file_type: str = "excel") -> Dict[str, Any]:
    """
    Importe un fichier (Excel, CSV ou JSON) dans un modèle Django spécifique.
    - model_name : ex. "masters.MasterEnrollment"
    - file : InMemoryUploadedFile ou chemin
    - file_type : "excel", "csv" ou "json"
    """

    try:
        model: Type = apps.get_model(model_name)
    except Exception as e:
        return {"ok": False, "error": f"Modèle introuvable : {e}"}

    # Lecture du fichier
    try:
        if file_type == "excel":
            df = pd.read_excel(file)
        elif file_type == "csv":
            df = pd.read_csv(file)
        elif file_type == "json":
            df = pd.DataFrame(json.load(file))
        else:
            return {"ok": False, "error": f"Type de fichier non supporté : {file_type}"}
    except Exception as e:
        return {"ok": False, "error": f"Erreur de lecture du fichier : {e}"}

    df = massage_dataframe(df)

    created, updated, errors = 0, 0, []

    # Boucle d’insertion
    for i, row in df.iterrows():
        try:
            obj_data = {field: row.get(field, None) for field in df.columns}
            obj = model(**obj_data)
            obj.full_clean()  # Validation Django (types, nullables)
            obj.save()
            created += 1
        except ValidationError as e:
            errors.append(f"Ligne {i+1}: {e.messages}")
        except Exception as e:
            errors.append(f"Ligne {i+1}: {e}")

    return {"ok": True, "created": created, "updated": updated, "errors": errors}


# ============================================================
# 📤 3️⃣ EXPORT GÉNÉRIQUE : Base Django vers fichier
# ============================================================

def export_to_excel(queryset, filename="export.xlsx") -> HttpResponse:
    """Exporte un queryset en fichier Excel (xlsx)."""
    df = pd.DataFrame(list(queryset.values()))
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False)
    buffer.seek(0)
    response = HttpResponse(buffer.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    response["Content-Disposition"] = f"attachment; filename={filename}"
    return response


def export_to_csv(queryset, filename="export.csv") -> HttpResponse:
    """Exporte un queryset en CSV."""
    df = pd.DataFrame(list(queryset.values()))
    buffer = io.StringIO()
    df.to_csv(buffer, index=False)
    response = HttpResponse(buffer.getvalue(), content_type="text/csv")
    response["Content-Disposition"] = f"attachment; filename={filename}"
    return response


def export_to_pdf(queryset, title="Rapport Export", filename="rapport.pdf") -> HttpResponse:
    """Exporte un queryset en PDF (liste tabulaire)."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]

    if queryset.exists():
        df = pd.DataFrame(list(queryset.values()))
        table_data = [list(df.columns)] + df.values.tolist()
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.cyan),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('INNERGRID', (0, 0), (-1, -1), 0.25, colors.grey),
            ('BOX', (0, 0), (-1, -1), 0.5, colors.black),
        ]))
        story.append(table)
    else:
        story.append(Paragraph("Aucune donnée disponible.", styles["Normal"]))

    doc.build(story)
    buffer.seek(0)
    response = HttpResponse(buffer.getvalue(), content_type="application/pdf")
    response["Content-Disposition"] = f"attachment; filename={filename}"
    return response


def export_to_word(queryset, title="Export Word", filename="rapport.docx") -> HttpResponse:
    """Exporte un queryset en Word (liste simple)."""
    doc = Document()
    doc.add_heading(title, 0)
    if queryset.exists():
        for obj in queryset:
            doc.add_paragraph(str(obj))
    else:
        doc.add_paragraph("Aucune donnée trouvée.")
    buffer = io.BytesIO()
    doc.save(buffer)
    response = HttpResponse(buffer.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    response["Content-Disposition"] = f"attachment; filename={filename}"
    return response


def export_to_pptx(queryset, title="Présentation ESFé", filename="rapport.pptx") -> HttpResponse:
    """Exporte un queryset en PowerPoint (tableaux simplifiés)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.text = f"Généré le {datetime.now().strftime('%d/%m/%Y %H:%M')}"

    for obj in queryset[:10]:
        p = text_frame.add_paragraph()
        p.text = str(obj)
        p.level = 1

    buffer = io.BytesIO()
    prs.save(buffer)
    response = HttpResponse(buffer.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response["Content-Disposition"] = f"attachment; filename={filename}"
    return response


# ============================================================
# 🧰 4️⃣ CLASSE GÉNÉRIQUE D’UTILISATION
# ============================================================

class ImportExportManager:
    """
    Classe utilitaire pour centraliser toutes les opérations d’import-export.
    Utilisable dans les vues Django ou les API REST.
    """

    @staticmethod
    def import_data(model_name: str, file, file_type="excel"):
        return import_from_file(model_name, file, file_type)

    @staticmethod
    def export_data(queryset, format="excel", title="Rapport ESFé", filename="export.xlsx"):
        format = format.lower()
        if format == "excel":
            return export_to_excel(queryset, filename)
        elif format == "csv":
            return export_to_csv(queryset, filename)
        elif format == "pdf":
            return export_to_pdf(queryset, title, filename)
        elif format == "word":
            return export_to_word(queryset, title, filename)
        elif format == "pptx":
            return export_to_pptx(queryset, title, filename)
        else:
            raise ValueError("Format d'export non supporté.")
